"""PowerPoint document loader for Microsoft PowerPoint files."""
from io import BytesIO
from typing import Any, BinaryIO

try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

from eva_rag.loaders.base import DocumentLoader, ExtractedDocument


class PowerPointLoader(DocumentLoader):
    """
    Loader for Microsoft PowerPoint presentations (.pptx).
    
    Features:
    - Extracts text from all slides
    - Includes slide titles, body text, and notes
    - Preserves slide order and structure
    - Extracts table data from slides
    - Includes speaker notes
    - Metadata: slide count, layout types
    
    Requires: python-pptx
    
    Example:
        >>> loader = PowerPointLoader(include_notes=True)
        >>> with open("quarterly_report.pptx", "rb") as f:
        ...     doc = loader.load_from_stream(f)
        >>> print(doc.metadata["slide_count"])  # 15
        >>> print(doc.text)  # All slides as markdown
    """
    
    def __init__(
        self,
        include_notes: bool = True,
        include_tables: bool = True,
    ):
        """
        Initialize PowerPoint loader.
        
        Args:
            include_notes: Whether to extract speaker notes
            include_tables: Whether to extract table data
        """
        if not PPTX_AVAILABLE:
            raise ImportError(
                "python-pptx is required for PowerPoint loading. "
                "Install with: pip install python-pptx"
            )
        
        self.include_notes = include_notes
        self.include_tables = include_tables
    
    def load(self, file: BinaryIO, filename: str) -> ExtractedDocument:
        """
        Load and extract text from PowerPoint document.
        
        Args:
            file: Binary file object
            filename: Original filename
            
        Returns:
            Extracted document with text and metadata
        """
        return self.load_from_stream(file)
    
    def load_from_stream(self, stream: BinaryIO) -> ExtractedDocument:
        """
        Load PowerPoint document from stream.
        
        Args:
            stream: Binary stream of PowerPoint file
            
        Returns:
            Extracted document with text and metadata
        """
        try:
            prs = Presentation(BytesIO(stream.read()))
        except Exception as e:
            return ExtractedDocument(
                text="",
                page_count=0,
                metadata={
                    "error": f"PowerPoint load error: {str(e)}",
                    "loader": "PowerPointLoader"
                }
            )
        
        # Extract slides
        slides_text = []
        slides_metadata = []
        
        for slide_idx, slide in enumerate(prs.slides, 1):
            slide_text, slide_meta = self._extract_slide(slide, slide_idx)
            slides_text.append(slide_text)
            slides_metadata.append(slide_meta)
        
        # Combine all slides
        text = "\n\n---\n\n".join(slides_text)
        
        # Extract core properties
        metadata = {
            "loader": "PowerPointLoader",
            "slide_count": len(prs.slides),
            "slides": slides_metadata,
        }
        
        if hasattr(prs, "core_properties") and prs.core_properties:
            props = prs.core_properties
            if props.title:
                metadata["title"] = props.title
            if props.author:
                metadata["author"] = props.author
            if props.created:
                metadata["created"] = str(props.created)
            if props.modified:
                metadata["modified"] = str(props.modified)
        
        return ExtractedDocument(
            text=text,
            page_count=len(prs.slides),
            metadata=metadata
        )
    
    def _extract_slide(self, slide, slide_num: int) -> tuple[str, dict[str, Any]]:
        """
        Extract text and metadata from a single slide.
        
        Args:
            slide: pptx slide object
            slide_num: Slide number (1-indexed)
            
        Returns:
            Tuple of (text, metadata)
        """
        lines = []
        lines.append(f"# Slide {slide_num}\n")
        
        # Extract text from shapes
        slide_text_parts = []
        table_count = 0
        
        for shape in slide.shapes:
            # Text from text frames
            if hasattr(shape, "text_frame"):
                text = shape.text_frame.text.strip()
                if text:
                    slide_text_parts.append(text)
            
            # Text from text boxes
            elif hasattr(shape, "text"):
                text = shape.text.strip()
                if text:
                    slide_text_parts.append(text)
            
            # Extract tables
            if self.include_tables and hasattr(shape, "table"):
                table_text = self._extract_table(shape.table)
                if table_text:
                    slide_text_parts.append(table_text)
                    table_count += 1
        
        # Add slide content
        if slide_text_parts:
            lines.append("\n".join(slide_text_parts))
        else:
            lines.append("*(Empty slide)*")
        
        # Extract speaker notes
        notes_text = None
        if self.include_notes and hasattr(slide, "notes_slide"):
            notes = slide.notes_slide
            if hasattr(notes, "notes_text_frame"):
                notes_text = notes.notes_text_frame.text.strip()
                if notes_text:
                    lines.append(f"\n**Speaker Notes:**\n{notes_text}")
        
        text = "\n".join(lines)
        
        metadata = {
            "slide_number": slide_num,
            "has_notes": bool(notes_text),
            "table_count": table_count,
            "shape_count": len(slide.shapes),
        }
        
        return text, metadata
    
    def _extract_table(self, table) -> str:
        """
        Extract text from a table.
        
        Args:
            table: pptx table object
            
        Returns:
            Markdown formatted table
        """
        if not table.rows or not table.columns:
            return ""
        
        lines = ["\n**Table:**\n"]
        
        # Extract table data
        rows_data = []
        for row in table.rows:
            row_data = []
            for cell in row.cells:
                text = cell.text.strip()
                row_data.append(text)
            rows_data.append(row_data)
        
        if not rows_data:
            return ""
        
        # Generate markdown table
        # Header (first row)
        header = "| " + " | ".join(rows_data[0]) + " |"
        separator = "| " + " | ".join(["---"] * len(rows_data[0])) + " |"
        lines.append(header)
        lines.append(separator)
        
        # Data rows
        for row in rows_data[1:]:
            line = "| " + " | ".join(row) + " |"
            lines.append(line)
        
        return "\n".join(lines)
