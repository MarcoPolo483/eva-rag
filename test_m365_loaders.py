"""Quick integration test for M365 loaders"""
import io
from src.eva_rag.loaders.csv_loader import CSVLoader
from src.eva_rag.loaders.excel_loader import ExcelLoader
from src.eva_rag.loaders.pptx_loader import PowerPointLoader

def test_csv():
    print('Testing CSVLoader...')
    csv_data = 'name,age,city\nAlice,30,Toronto\nBob,25,Vancouver\nCharlie,35,Montreal'
    csv_stream = io.BytesIO(csv_data.encode('utf-8'))
    csv_loader = CSVLoader()
    csv_doc = csv_loader.load_from_stream(csv_stream)
    print(f'✅ CSV: {csv_doc.metadata["row_count"]} rows, {csv_doc.metadata["column_count"]} columns')
    print(f'   Columns: {", ".join(csv_doc.metadata["columns"])}')
    return csv_doc

def test_excel():
    print('\nTesting ExcelLoader...')
    from openpyxl import Workbook
    wb = Workbook()
    ws = wb.active
    ws.title = 'Employment'
    ws.append(['Province', 'Unemployment Rate', 'Year'])
    ws.append(['Ontario', 5.2, 2024])
    ws.append(['Quebec', 4.8, 2024])
    excel_stream = io.BytesIO()
    wb.save(excel_stream)
    excel_stream.seek(0)
    excel_loader = ExcelLoader()
    excel_doc = excel_loader.load_from_stream(excel_stream)
    print(f'✅ Excel: {excel_doc.metadata["sheet_count"]} sheet(s), {len(excel_doc.text)} chars')
    print(f'   Sheets: {", ".join(excel_doc.metadata["sheet_names"])}')
    return excel_doc

def test_powerpoint():
    print('\nTesting PowerPointLoader...')
    from pptx import Presentation
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = 'Employment Trends'
    slide.placeholders[1].text = 'Canadian labor market statistics show positive growth.'
    pptx_stream = io.BytesIO()
    prs.save(pptx_stream)
    pptx_stream.seek(0)
    pptx_loader = PowerPointLoader()
    pptx_doc = pptx_loader.load_from_stream(pptx_stream)
    print(f'✅ PowerPoint: {pptx_doc.metadata["slide_count"]} slide(s), {len(pptx_doc.text)} chars')
    return pptx_doc

if __name__ == '__main__':
    print('=== Testing M365 Loaders ===\n')
    
    csv_doc = test_csv()
    excel_doc = test_excel()
    pptx_doc = test_powerpoint()
    
    print('\n=== Sample Output ===')
    print('\nCSV Text Preview:')
    print(csv_doc.text[:200])
    
    print('\n\nExcel Text Preview:')
    print(excel_doc.text[:200])
    
    print('\n\nPowerPoint Text Preview:')
    print(pptx_doc.text[:200])
    
    print('\n\n✅ All M365 loaders working correctly!')
